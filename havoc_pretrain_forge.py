import os
import json
import random
import re

# Third-party libs - make sure these are installed:
# pip install pandas pyarrow lxml pypdf python-docx python-pptx

import pandas as pd
from lxml import etree
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# WHERE THE DATA LIVES
DATA_FOLDER = r"C:\Users\ScottT\Desktop\HAVOC_DATA"

# WHERE TO WRITE PRETRAIN OUTPUT
OUTPUT_DIR = r"C:\Users\ScottT\Desktop\HAVOC_PRETRAIN"

# Validation fraction (2% of samples)
VAL_FRACTION = 0.02

os.makedirs(OUTPUT_DIR, exist_ok=True)
train_path = os.path.join(OUTPUT_DIR, "pretrain_train.jsonl")
val_path = os.path.join(OUTPUT_DIR, "pretrain_val.jsonl")

random.seed(1337)


def clean_text(s: str) -> str:
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = re.sub(r'\n{3,}', '\n\n', s)
    s = re.sub(r'[ \t]{2,}', ' ', s)
    return s.strip()


def emit(text: str, train_fp, val_fp):
    text = clean_text(text)
    if len(text) < 50:  # drop ultra-short junk
        return
    rec = {"text": text}
    line = json.dumps(rec, ensure_ascii=False)
    if random.random() < VAL_FRACTION:
        val_fp.write(line + "\n")
    else:
        train_fp.write(line + "\n")


def process_parquet(path, train_fp, val_fp):
    try:
        df = pd.read_parquet(path)
    except Exception as e:
        print(f"[PARQUET] Failed {path}: {e}")
        return
    text_cols = [c for c in df.columns if df[c].dtype == object]
    if not text_cols:
        return
    print(f"[PARQUET] {path} -> columns {text_cols}")
    for col in text_cols:
        for v in df[col].dropna():
            if isinstance(v, str):
                emit(v, train_fp, val_fp)


def process_xml(path, train_fp, val_fp):
    try:
        tree = etree.parse(path)
        root = tree.getroot()
        text = " ".join(root.itertext())
        emit(text, train_fp, val_fp)
        print(f"[XML] {path}")
    except Exception as e:
        print(f"[XML] Failed {path}: {e}")


def process_pdf(path, train_fp, val_fp):
    try:
        reader = PdfReader(path)
    except Exception as e:
        print(f"[PDF] Failed open {path}: {e}")
        return
    texts = []
    for page in reader.pages:
        try:
            t = page.extract_text()
        except Exception:
            t = None
        if t:
            texts.append(t)
    if texts:
        emit("\n\n".join(texts), train_fp, val_fp)
        print(f"[PDF] {path}")


def process_docx(path, train_fp, val_fp):
    try:
        doc = DocxDocument(path)
    except Exception as e:
        print(f"[DOCX] Failed {path}: {e}")
        return
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    if paras:
        emit("\n\n".join(paras), train_fp, val_fp)
        print(f"[DOCX] {path}")


def process_pptx(path, train_fp, val_fp):
    try:
        pres = Presentation(path)
    except Exception as e:
        print(f"[PPTX] Failed {path}: {e}")
        return
    texts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    if texts:
        emit("\n\n".join(texts), train_fp, val_fp)
        print(f"[PPTX] {path}")


def extract_strings(obj):
    out = []
    if isinstance(obj, str):
        out.append(obj)
    elif isinstance(obj, dict):
        for v in obj.values():
            out.extend(extract_strings(v))
    elif isinstance(obj, (list, tuple)):
        for v in obj:
            out.extend(extract_strings(v))
    return out


def process_json(path, train_fp, val_fp):
    try:
        with open(path, "r", encoding="utf8") as f:
            data = json.load(f)
    except Exception as e:
        print(f"[JSON] Failed {path}: {e}")
        return
    texts = extract_strings(data)
    if not texts:
        return
    emit("\n\n".join(texts), train_fp, val_fp)
    print(f"[JSON] {path} ({len(texts)} strings)")


def process_jsonl(path, train_fp, val_fp):
    print(f"[JSONL] {path}")
    with open(path, "r", encoding="utf8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                obj = json.loads(line)
            except Exception:
                continue
            text = None
            # common keys first
            for key in ("text", "content", "body"):
                if isinstance(obj, dict) and key in obj and isinstance(obj[key], str):
                    text = obj[key]
                    break
            if text is None:
                strings = extract_strings(obj)
                if strings:
                    text = "\n\n".join(strings)
            if text:
                emit(text, train_fp, val_fp)


def process_txt(path, train_fp, val_fp):
    try:
        with open(path, "r", encoding="utf8", errors="ignore") as f:
            text = f.read()
    except Exception as e:
        print(f"[TXT] Failed {path}: {e}")
        return
    emit(text, train_fp, val_fp)
    print(f"[TXT] {path}")


def main():
    with open(train_path, "w", encoding="utf8") as train_fp, \
         open(val_path, "w", encoding="utf8") as val_fp:

        for root, dirs, files in os.walk(DATA_FOLDER):
            for name in files:
                path = os.path.join(root, name)
                ext = os.path.splitext(name)[1].lower()

                try:
                    if ext == ".parquet":
                        process_parquet(path, train_fp, val_fp)
                    elif ext == ".xml":
                        process_xml(path, train_fp, val_fp)
                    elif ext == ".pdf":
                        process_pdf(path, train_fp, val_fp)
                    elif ext in (".doc", ".docx"):
                        process_docx(path, train_fp, val_fp)
                    elif ext in (".ppt", ".pptx"):
                        process_pptx(path, train_fp, val_fp)
                    elif ext == ".json":
                        process_json(path, train_fp, val_fp)
                    elif ext == ".jsonl":
                        process_jsonl(path, train_fp, val_fp)
                    elif ext in (".txt", ".csv", ".html"):
                        process_txt(path, train_fp, val_fp)
                    else:
                        # Skipping non-text / archive / binary types:
                        # .7z, .zip, .rsmi, .lmdb, .numbers, .pages, .xyz, etc.
                        continue
                except Exception as e:
                    print(f"[ERROR] {path}: {e}")

    print("DONE. Pretrain dataset written to:")
    print("  ", train_path)
    print("  ", val_path)


if __name__ == "__main__":
    main()
